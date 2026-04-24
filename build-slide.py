#!/usr/bin/env python3
"""Build the 'How PMM Uses AI Tools' PowerPoint slide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──
HF_GREEN = RGBColor(0x1B, 0x8A, 0x4A)
HF_DARK_GREEN = RGBColor(0x0B, 0x5C, 0x30)
HF_LIME = RGBColor(0xC8, 0xF0, 0x65)
HF_CREAM = RGBColor(0xFF, 0xFF, 0xFF)
HF_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HF_DARK = RGBColor(0x1A, 0x1A, 0x1A)
HF_MUTED = RGBColor(0x6B, 0x6B, 0x6B)
GEMINI_BLUE = RGBColor(0x42, 0x85, 0xF4)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xED)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFD)
LIGHT_GREY_BG = RGBColor(0xF0, 0xF0, 0xF0)
CARD_BORDER = RGBColor(0xE8, 0xE8, 0xE8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# Set background
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = HF_WHITE


# ── Helper functions ──
def add_shape(left, top, width, height, fill_color=None, border_color=None, border_width=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or HF_WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        # Adjustment value for corner radius (0-50000)
        shape.adjustments[0] = radius
    return shape

def add_rect(left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or HF_WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(left, top, width, height, text, font_size=12, bold=False, color=HF_DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(left, top, width, height, lines, font_name='Calibri'):
    """lines: list of (text, font_size, bold, color, spacing_after) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, bold, color, spacing) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing) if spacing else Pt(0)
    return txBox


# ══════════════════════════════════════════
# HEADER AREA
# ══════════════════════════════════════════

# Topic badge
badge = add_shape(Inches(0.6), Inches(0.4), Inches(2.0), Inches(0.32), fill_color=LIGHT_GREEN_BG, radius=0.5)
add_text(Inches(0.6), Inches(0.4), Inches(2.0), Inches(0.32),
         "PRODUCT MARKETING", font_size=8, bold=True, color=HF_GREEN, alignment=PP_ALIGN.CENTER)

# Title
title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.82), Inches(7), Inches(0.7))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run1 = p.add_run()
run1.text = "How PMM Uses "
run1.font.size = Pt(36)
run1.font.bold = True
run1.font.color.rgb = HF_DARK
run1.font.name = 'Calibri'
run2 = p.add_run()
run2.text = "AI Tools"
run2.font.size = Pt(36)
run2.font.bold = True
run2.font.color.rgb = HF_GREEN
run2.font.name = 'Calibri'

# Subtitle
add_text(Inches(0.6), Inches(1.52), Inches(6.5), Inches(0.6),
         "Turning a campaign idea into an interactive, shareable prototype — using Claude Code for strategy & development, and Gemini for visual asset generation.",
         font_size=11, color=HF_MUTED)

# Presenter info (top right)
add_text(Inches(10.5), Inches(0.5), Inches(2.5), Inches(0.28),
         "Kieran Enright", font_size=13, bold=True, color=HF_DARK, alignment=PP_ALIGN.RIGHT)
add_text(Inches(10.5), Inches(0.78), Inches(2.5), Inches(0.24),
         "Product Marketing Manager", font_size=10, color=HF_MUTED, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════
# LEFT SIDE: WORKFLOW STEPS
# ══════════════════════════════════════════

left_start = Inches(0.6)
top_start = Inches(2.35)
step_height = Inches(1.15)

steps = [
    ("1", "💬", "Describe the campaign idea",
     "Plain English conversation — no formal brief, no templates. Just describe the concept, audience, and goals.",
     HF_DARK, None, None),
    ("2", "🧠", "AI structures the strategy",
     "Claude Code builds the full messaging framework, channel strategy, timeline, and validation plan.",
     HF_GREEN, "Claude Code", "claude"),
    ("3", "🎨", "Generate visual assets",
     "Gemini creates campaign mockups — social posts, phone screens, visual concepts — from text-based art direction.",
     GEMINI_BLUE, "Gemini", "gemini"),
    ("4", "🚀", "Build & deploy the prototype",
     "Claude Code writes full interactive HTML — scroll animations, carousels, lightboxes — and ships to a live URL.",
     HF_GREEN, "Claude Code + GitHub Pages", "both"),
]

for i, (num, emoji, title, desc, accent, tool_label, tool_type) in enumerate(steps):
    y = top_start + Emu(int(step_height * i))

    # Step number circle
    circle = add_shape(left_start, y, Inches(0.42), Inches(0.42), fill_color=accent, radius=0.5)
    # Number/emoji on circle
    add_text(left_start, y, Inches(0.42), Inches(0.42), emoji,
             font_size=16, alignment=PP_ALIGN.CENTER, color=HF_WHITE)

    # Connector line (except last step)
    if i < len(steps) - 1:
        line_x = left_start + Inches(0.19)
        line_y = y + Inches(0.44)
        line = add_rect(line_x, line_y, Inches(0.04), Inches(0.68), fill_color=RGBColor(0xDD, 0xDD, 0xDD))

    # Step label
    add_text(left_start + Inches(0.6), y - Inches(0.02), Inches(1.2), Inches(0.2),
             f"STEP {num}", font_size=8, bold=True, color=HF_MUTED)

    # Step title
    add_text(left_start + Inches(0.6), y + Inches(0.16), Inches(5), Inches(0.28),
             title, font_size=14, bold=True, color=HF_DARK)

    # Step description
    add_text(left_start + Inches(0.6), y + Inches(0.42), Inches(5.2), Inches(0.4),
             desc, font_size=9.5, color=HF_MUTED)

    # Tool badge(s)
    if tool_label:
        badge_y = y + Inches(0.8)
        if tool_type == "claude":
            b = add_shape(left_start + Inches(0.6), badge_y, Inches(1.2), Inches(0.24),
                         fill_color=LIGHT_GREEN_BG, radius=0.4)
            add_text(left_start + Inches(0.6), badge_y, Inches(1.2), Inches(0.24),
                    "Claude Code", font_size=8, bold=True, color=HF_GREEN, alignment=PP_ALIGN.CENTER)
        elif tool_type == "gemini":
            b = add_shape(left_start + Inches(0.6), badge_y, Inches(0.9), Inches(0.24),
                         fill_color=LIGHT_BLUE_BG, radius=0.4)
            add_text(left_start + Inches(0.6), badge_y, Inches(0.9), Inches(0.24),
                    "Gemini", font_size=8, bold=True, color=GEMINI_BLUE, alignment=PP_ALIGN.CENTER)
        elif tool_type == "both":
            b1 = add_shape(left_start + Inches(0.6), badge_y, Inches(1.2), Inches(0.24),
                          fill_color=LIGHT_GREEN_BG, radius=0.4)
            add_text(left_start + Inches(0.6), badge_y, Inches(1.2), Inches(0.24),
                    "Claude Code", font_size=8, bold=True, color=HF_GREEN, alignment=PP_ALIGN.CENTER)
            b2 = add_shape(left_start + Inches(1.9), badge_y, Inches(1.3), Inches(0.24),
                          fill_color=LIGHT_GREY_BG, radius=0.4)
            add_text(left_start + Inches(1.9), badge_y, Inches(1.3), Inches(0.24),
                    "GitHub Pages", font_size=8, bold=True, color=HF_DARK, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# RIGHT SIDE: BROWSER MOCKUP
# ══════════════════════════════════════════

browser_left = Inches(7.0)
browser_top = Inches(2.35)
browser_width = Inches(5.8)

# Browser frame outer
browser_frame = add_shape(browser_left, browser_top, browser_width, Inches(3.5),
                          fill_color=HF_WHITE, border_color=CARD_BORDER, radius=0.03)

# Browser bar background
browser_bar = add_rect(browser_left + Inches(0.01), browser_top + Inches(0.01),
                       browser_width - Inches(0.02), Inches(0.4), fill_color=RGBColor(0xFA, 0xFA, 0xFA))

# Traffic light dots
dot_y = browser_top + Inches(0.15)
add_circle(browser_left + Inches(0.2), dot_y, Inches(0.1), RGBColor(0xFF, 0x5F, 0x57))
add_circle(browser_left + Inches(0.38), dot_y, Inches(0.1), RGBColor(0xFE, 0xBC, 0x2E))
add_circle(browser_left + Inches(0.56), dot_y, Inches(0.1), RGBColor(0x28, 0xC8, 0x40))

# URL bar
url_bar = add_shape(browser_left + Inches(0.8), browser_top + Inches(0.1),
                    Inches(4.7), Inches(0.24), fill_color=HF_WHITE,
                    border_color=RGBColor(0xE0, 0xE0, 0xE0), radius=0.3)
add_text(browser_left + Inches(0.9), browser_top + Inches(0.1), Inches(4.5), Inches(0.24),
         "kieranenright.github.io/marketing-prototypes/cookbook-campaign-strategy",
         font_size=7.5, color=HF_MUTED)

# Browser content area (green gradient simulation)
content_area = add_rect(browser_left + Inches(0.01), browser_top + Inches(0.42),
                        browser_width - Inches(0.02), Inches(3.06), fill_color=HF_GREEN)

# Decorative lighter circle (top right of content)
deco_circle = add_circle(browser_left + Inches(3.8), browser_top + Inches(0.2),
                         Inches(2.8), RGBColor(0x22, 0x96, 0x55))
deco_circle.fill.solid()
deco_circle.fill.fore_color.rgb = RGBColor(0x22, 0x96, 0x55)

# Proto badge
proto_badge = add_shape(browser_left + Inches(1.6), browser_top + Inches(1.2),
                        Inches(2.6), Inches(0.26),
                        fill_color=RGBColor(0x1F, 0x7A, 0x47), radius=0.5)
add_text(browser_left + Inches(1.6), browser_top + Inches(1.2), Inches(2.6), Inches(0.26),
         "Q2 2026 CAMPAIGN STRATEGY", font_size=7, bold=True,
         color=HF_LIME, alignment=PP_ALIGN.CENTER)

# Proto title line 1
add_text(browser_left + Inches(0.4), browser_top + Inches(1.6), Inches(5.0), Inches(0.4),
         "Don't just save it.", font_size=24, bold=True, color=HF_WHITE, alignment=PP_ALIGN.CENTER)

# Proto title line 2 (lime)
add_text(browser_left + Inches(0.4), browser_top + Inches(2.0), Inches(5.0), Inches(0.4),
         "Cook it.", font_size=24, bold=True, color=HF_LIME, alignment=PP_ALIGN.CENTER)

# Proto subtitle
add_text(browser_left + Inches(0.4), browser_top + Inches(2.45), Inches(5.0), Inches(0.3),
         "Cookbook — Efficient Inspired Campaign", font_size=9,
         color=RGBColor(0x8B, 0xBF, 0x9E), alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# STAT CARDS
# ══════════════════════════════════════════

stats_top = Inches(6.05)
card_width = Inches(1.82)
card_height = Inches(0.85)
card_gap = Inches(0.12)

stats_data = [
    ("0", "designers\nrequired"),
    ("0", "developers\nrequired"),
    ("1", "marketer +\nAI tools"),
]

for i, (number, label) in enumerate(stats_data):
    x = browser_left + Emu(int((card_width + card_gap) * i))

    card = add_shape(x, stats_top, card_width, card_height,
                     fill_color=HF_WHITE, border_color=CARD_BORDER, radius=0.08)

    add_text(x, stats_top + Inches(0.08), card_width, Inches(0.38),
             number, font_size=28, bold=True, color=HF_GREEN, alignment=PP_ALIGN.CENTER)

    # Label (multiline)
    label_box = slide.shapes.add_textbox(x, stats_top + Inches(0.48), card_width, Inches(0.35))
    tf = label_box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(label.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(8)
        p.font.color.rgb = HF_MUTED
        p.font.bold = True
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)


# ══════════════════════════════════════════
# FOOTER
# ══════════════════════════════════════════

footer_y = Inches(7.1)

add_text(Inches(0.6), footer_y, Inches(3), Inches(0.24),
         "Brand-Wide Meeting — Q2 2026", font_size=9, color=HF_MUTED)

# Tool legend (right side)
tools_legend = [
    ("Claude Code", HF_GREEN),
    ("Gemini", GEMINI_BLUE),
    ("GitHub Pages", HF_DARK),
]

legend_x = Inches(9.5)
for tool_name, dot_color in tools_legend:
    dot = add_circle(legend_x, footer_y + Inches(0.06), Inches(0.1), dot_color)
    add_text(legend_x + Inches(0.16), footer_y, Inches(1.2), Inches(0.22),
             tool_name, font_size=9, bold=True, color=HF_MUTED)
    legend_x += Inches(1.35)


# ══════════════════════════════════════════
# SPEAKER NOTES
# ══════════════════════════════════════════

notes_slide = slide.notes_slide
notes_tf = notes_slide.notes_text_frame
notes_tf.text = ""

notes_lines = [
    "THE PROBLEM: Communicating campaign strategy through static slides and docs doesn't convey the experience. Stakeholders skim decks — they engage with prototypes.",
    "",
    "THE WORKFLOW: I described the Cookbook campaign concept to Claude Code in plain conversation. It structured the full strategy — messaging angles, channel mix, 13-week timeline, validation roadmap — then built the entire interactive prototype as a single HTML file.",
    "",
    "GEMINI'S ROLE: I used Gemini to generate all the visual campaign mockups — social posts, phone screens, creative concepts — using text prompts based on our brand guidelines. No designer needed at the concept stage.",
    "",
    "THE RESULT: A fully interactive, scroll-driven strategy document with animations, carousels, lightbox viewers, and a colour-coded deployment calendar. Anyone with the link can experience it — no special tools needed.",
    "",
    "WHAT THIS MEANS FOR PMM: What used to require cross-functional coordination (strategist + designer + developer) can now be prototyped by a single marketer. The cost of iteration is a conversation, not a ticket.",
    "",
    "KEY STAT: The prototype maps out ~290 content pieces across organic, paid, and ASO channels with a full validation timeline.",
]

for i, line in enumerate(notes_lines):
    if i == 0:
        p = notes_tf.paragraphs[0]
    else:
        p = notes_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.name = 'Calibri'


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "prototypes", "ai-workflow-slide.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
